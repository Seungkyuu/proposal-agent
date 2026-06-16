import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import os

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

slide_w = prs.slide_width
slide_h = prs.slide_height
print(f"Slides: {len(prs.slides)}")
print(f"Slide size: {slide_w.inches:.2f}\" x {slide_h.inches:.2f}\"")
print()

for i, slide in enumerate(prs.slides, 1):
    print(f"=== Slide {i} ===")
    for shape in slide.shapes:
        l = shape.left or 0
        t = shape.top or 0
        w = shape.width or 0
        h = shape.height or 0
        r = l + w
        b = t + h
        overflow_right = r > slide_w
        overflow_bottom = b > slide_h
        neg_pos = l < 0 or t < 0

        info = f"  Shape '{shape.name}' | pos=({Emu(l).inches:.2f}\", {Emu(t).inches:.2f}\") size=({Emu(w).inches:.2f}\"x{Emu(h).inches:.2f}\")"
        flags = []
        if overflow_right:
            flags.append(f"OVERFLOW_RIGHT by {Emu(r - slide_w).inches:.3f}\"")
        if overflow_bottom:
            flags.append(f"OVERFLOW_BOTTOM by {Emu(b - slide_h).inches:.3f}\"")
        if neg_pos:
            flags.append("NEGATIVE_POSITION")

        if shape.has_text_frame:
            tf = shape.text_frame
            text_preview = tf.text[:80].replace('\n', ' | ')
            info += f"\n    Text: '{text_preview}'"
            # Check each paragraph's runs for font size
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pass  # sizes ok

        print(info)
        if flags:
            for f in flags:
                print(f"    *** {f} ***")
    print()
