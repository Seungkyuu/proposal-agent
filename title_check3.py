import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu, Pt

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

# Font sizes in python-pptx are stored as EMU (1 pt = 12700 EMU)
# So 381000 EMU = 381000/12700 = 30pt
def emu_to_pt(emu_val):
    if emu_val is None:
        return None
    return emu_val / 12700.0

print("=== Title shapes with correct point sizes ===")
print()

for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t_in = Emu(shape.top or 0).inches
        h_in = Emu(shape.height or 0).inches
        if t_in > 0.6:
            continue
        tf = shape.text_frame
        text = tf.text
        print(f"Slide {i}: top={t_in:.3f}\" height={h_in:.3f}\"")
        for pi, para in enumerate(tf.paragraphs):
            ptext = para.text[:60]
            pfont = None
            # Try run font first
            for run in para.runs:
                if run.font.size:
                    pfont = emu_to_pt(run.font.size)
                    break
            # Try para default font
            if pfont is None:
                try:
                    pfont = emu_to_pt(para._pPr.attrib.get('{http://schemas.openxmlformats.org/drawingml/2006/main}sz'))
                except:
                    pass
            print(f"  para[{pi}]: '{ptext}' font={pfont}pt")
            if pfont:
                # At 30pt, line height approx 30 * 1.2 / 72 = 0.500"
                line_h = pfont * 1.2 / 72
                print(f"    -> line height ~{line_h:.3f}\" | box height {h_in:.3f}\"")
        # Estimate if two lines fit
        max_pt = 0
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    pt = emu_to_pt(run.font.size)
                    if pt > max_pt:
                        max_pt = pt
        if max_pt > 0:
            line_h = max_pt * 1.2 / 72
            needed = line_h * 2 + 0.05
            status = "OK" if needed <= h_in else "*** CLIPPING RISK ***"
            print(f"  => two-line need: {needed:.3f}\" vs box {h_in:.3f}\" => {status}")
        print()

print()
print("=== Slide 7 special check (smaller font) ===")
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if not shape.has_text_frame:
        continue
    t_in = Emu(shape.top or 0).inches
    h_in = Emu(shape.height or 0).inches
    if t_in > 0.5:
        continue
    print(f"  height={h_in:.3f}\"")
    for pi, para in enumerate(shape.text_frame.paragraphs):
        for run in para.runs:
            if run.font.size:
                pt = emu_to_pt(run.font.size)
                line_h = pt * 1.2 / 72
                needed = line_h * 2 + 0.05
                print(f"  font={pt:.1f}pt line={line_h:.3f}\" 2-line={needed:.3f}\" box={h_in:.3f}\"")
                status = "OK" if needed <= h_in else "*** CLIPPING RISK ***"
                print(f"  => {status}")
                break
