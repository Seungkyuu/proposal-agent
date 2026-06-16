import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu, Pt

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

SLIDE_H_IN = Emu(prs.slide_height).inches

print("=== Title Shape Analysis (Slides 3-9) ===")
print("Checking: title box height 0.98\" — is it tall enough for two lines?")
print()

for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = Emu(shape.top or 0).inches
        h = Emu(shape.height or 0).inches
        # Title boxes: near top, have pipe separator
        if t < 0.5 and '|' in shape.text_frame.text:
            tf = shape.text_frame
            text = tf.text
            parts = [p.strip() for p in text.split('\n')]
            print(f"Slide {i} title:")
            print(f"  box: top={t:.2f}\" height={h:.2f}\" bottom={(t+h):.2f}\"")
            # Get font sizes
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt_size = Pt(run.font.size).pt
                        # Rough line height: font_size * 1.2 (in points), convert to inches
                        # 1 inch = 72 points
                        line_h_in = pt_size * 1.2 / 72
                        print(f"  font: {pt_size:.0f}pt, line height ~{line_h_in:.3f}\"")
                        # For two lines:
                        two_line_h = line_h_in * 2 + 0.05  # small padding
                        print(f"  two-line estimate: {two_line_h:.3f}\" vs box height {h:.2f}\"")
                        if two_line_h > h:
                            print(f"  *** CLIPPING RISK: box {h:.3f}\" < two-line height {two_line_h:.3f}\" ***")
                        else:
                            print(f"  OK: box has {h - two_line_h:.3f}\" spare for two lines")

            # Check if title text has pipe indicating two separate visual lines
            raw_parts = [p.strip() for p in text.split('|')]
            print(f"  pipe-parts: {len(raw_parts)}")
            for j, p in enumerate(raw_parts):
                print(f"    [{j+1}] '{p[:60]}'")
            print()

print()
print("=== Slide 1 Stacking Check ===")
slide1 = prs.slides[0]
shapes_s1 = sorted(slide1.shapes, key=lambda s: Emu(s.top or 0).inches)
for s in shapes_s1:
    t = Emu(s.top or 0).inches
    h = Emu(s.height or 0).inches
    b = t + h
    text = s.text_frame.text[:60] if s.has_text_frame else "(no text)"
    print(f"  top={t:.2f}\" bottom={b:.2f}\" | '{text}'")

print()
print("=== Slide 7 title height (smaller font) ===")
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        t = Emu(shape.top or 0).inches
        h = Emu(shape.height or 0).inches
        if t < 0.5:
            print(f"  top={t:.2f}\" height={h:.2f}\" text='{shape.text_frame.text[:60]}'")
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        pt_size = Pt(run.font.size).pt
                        line_h_in = pt_size * 1.2 / 72
                        two_line_h = line_h_in * 2 + 0.05
                        print(f"  font={pt_size:.0f}pt line_h={line_h_in:.3f}\" 2-line={two_line_h:.3f}\" box_h={h:.3f}\"")
