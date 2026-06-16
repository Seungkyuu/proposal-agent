import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu, Pt
from lxml import etree

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def emu_to_in(e):
    return Emu(e or 0).inches

print("=== Table analysis on each slide ===")
for i, slide in enumerate(prs.slides, 1):
    tables = [s for s in slide.shapes if s.shape_type == 19]  # MSO_SHAPE_TYPE.TABLE
    if not tables:
        continue
    for tbl_shape in tables:
        t_in = emu_to_in(tbl_shape.top)
        h_in = emu_to_in(tbl_shape.height)
        w_in = emu_to_in(tbl_shape.width)
        b_in = t_in + h_in
        print(f"Slide {i}: Table '{tbl_shape.name}' top={t_in:.3f}\" height={h_in:.3f}\" bottom={b_in:.3f}\"")
        tbl = tbl_shape.table
        print(f"  rows={len(tbl.rows)} cols={len(tbl.columns)}")
        for ri, row in enumerate(tbl.rows):
            row_h = emu_to_in(row.height)
            row_texts = []
            for cell in row.cells:
                row_texts.append(cell.text[:30])
            print(f"  row[{ri}] h={row_h:.3f}\" | {row_texts}")
        print()

print()
print("=== Section header text boxes (checking for background colors) ===")
for i, slide in enumerate(prs.slides, 1):
    if i < 3:
        continue
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t_in = emu_to_in(shape.top)
        h_in = emu_to_in(shape.height)
        # Section headers typically 1.1-1.5" from top
        if 1.0 <= t_in <= 1.5:
            text = shape.text_frame.text[:50]
            w_in = emu_to_in(shape.width)
            print(f"  Slide {i}: '{text}' top={t_in:.3f}\" h={h_in:.3f}\" w={w_in:.2f}\"")
            # Check fill
            sp = shape._element
            spPr = sp.find('.//{%s}spPr' % NS) or sp.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
            if spPr is not None:
                solidFill = spPr.find('.//{%s}solidFill' % NS)
                if solidFill is not None:
                    srgbClr = solidFill.find('{%s}srgbClr' % NS)
                    if srgbClr is not None:
                        print(f"    fill color: #{srgbClr.get('val', 'unknown')}")
