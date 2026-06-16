"""
pptx_builder.py  --  alpaco proposal text-box scaffold renderer
==============================================================
Method: No template cloning. TextBox objects placed directly on blank slides.
        Coordinates/sizes/styles loaded from templates/layout_positions.json.

Simple style:
  Background:       #FFFFFF
  Main text:        #1A1A2E
  Secondary text:   #6B7280
  Accent (headers): client brand HEX (from proposal JSON)
  Accent text:      #FFFFFF
  Font:             Malgun Gothic (Malgun Gothic)

Usage:
  pip install python-pptx
  python renderer/pptx_builder.py <proposal_json> <output_pptx>
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# -- Constants -----------------------------------------------------------------
SLIDE_W_CM = 33.867
SLIDE_H_CM = 19.05
FONT_FAMILY = "Malgun Gothic"

COLOR_BG    = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_SUB   = RGBColor(0x6B, 0x72, 0x80)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

POSITIONS_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "templates", "layout_positions.json"
)


# -- Color helpers -------------------------------------------------------------

def hex_to_rgb(hex_str):
    """'#RRGGBB' -> RGBColor"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def resolve_color(color_str, accent_color):
    """
    color_str: '#RRGGBB' | 'ACCENT' | None
    returns:   RGBColor | None
    """
    if color_str is None:
        return None
    if str(color_str).upper() == "ACCENT":
        return accent_color
    return hex_to_rgb(color_str)


# -- Slide helpers -------------------------------------------------------------

def add_blank_slide(prs):
    """Add a completely blank slide (layout index 6) with white background."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG
    return slide


def add_textbox(slide, box_def, accent_color, text, bullet=False):
    """
    Add a styled TextBox to the slide.

    box_def : dict from layout_positions.json
    text    : string, use \\n for line breaks
    bullet  : prepend bullet char to each non-empty line
    """
    left   = Cm(box_def["left_cm"])
    top    = Cm(box_def["top_cm"])
    width  = Cm(box_def["width_cm"])
    height = Cm(box_def["height_cm"])

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    font_size    = Pt(box_def["font_size_pt"])
    bold         = box_def.get("bold", False)
    color_str    = box_def.get("color", "#1A1A2E")
    align_str    = box_def.get("align", "left")
    bg_color_str = box_def.get("bg_color")

    font_color = resolve_color(color_str, accent_color)

    # Background fill for header boxes
    if bg_color_str:
        bg_rgb = resolve_color(bg_color_str, accent_color)
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_rgb
        tf.margin_left   = Cm(0.25)
        tf.margin_right  = Cm(0.25)
        tf.margin_top    = Cm(0.1)
        tf.margin_bottom = Cm(0.1)

    align_map = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }
    pptx_align = align_map.get(align_str, PP_ALIGN.LEFT)

    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = pptx_align
        display = ("- " + line) if (bullet and line.strip()) else line
        run = para.add_run()
        run.text = display
        run.font.name = FONT_FAMILY
        run.font.size = font_size
        run.font.bold = bold
        if font_color:
            run.font.color.rgb = font_color

    return txBox


def add_table(slide, table_def, accent_color, headers, rows):
    """
    Add a styled table to the slide.

    table_def : dict from layout_positions.json "table" key
    headers   : list of column header strings
    rows      : list of row lists
    """
    from lxml import etree

    left   = Cm(table_def["left_cm"])
    top    = Cm(table_def["top_cm"])
    width  = Cm(table_def["width_cm"])
    height = Cm(table_def["height_cm"])

    n_rows = 1 + len(rows)
    n_cols = len(headers)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    header_bg    = resolve_color(table_def.get("header_bg", "ACCENT"), accent_color)
    header_pt    = table_def.get("header_font_pt", 12)
    header_bold  = table_def.get("header_bold", True)
    header_color = resolve_color(table_def.get("header_color", "#FFFFFF"), accent_color)
    data_pt      = table_def.get("data_font_pt", 11)
    data_bold    = table_def.get("data_bold", False)
    data_color   = resolve_color(table_def.get("data_color", "#1A1A2E"), accent_color)
    alt_bg_str   = table_def.get("alt_row_bg")
    alt_bg       = hex_to_rgb(alt_bg_str) if alt_bg_str else None

    # Column widths
    col_widths_pct = table_def.get("col_widths_pct")
    if col_widths_pct and len(col_widths_pct) == n_cols:
        total_emu = int(width)
        for ci, pct in enumerate(col_widths_pct):
            tbl.columns[ci].width = int(total_emu * pct / 100)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = str(hdr)
        _style_cell(cell, FONT_FAMILY, header_pt, header_bold, header_color, header_bg)

    # Data rows
    for ri, row_data in enumerate(rows):
        bg = alt_bg if (alt_bg and ri % 2 == 1) else None
        for ci in range(min(len(row_data), n_cols)):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(row_data[ci]) if row_data[ci] is not None else ""
            _style_cell(cell, FONT_FAMILY, data_pt, data_bold, data_color, bg)

    return tbl_shape


def _style_cell(cell, font_name, font_pt, bold, font_color, bg_color):
    """Apply font and background to a table cell."""
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name  = font_name
            run.font.size  = Pt(font_pt)
            run.font.bold  = bold
            if font_color:
                run.font.color.rgb = font_color

    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        old  = tcPr.find(qn("a:solidFill"))
        if old is not None:
            tcPr.remove(old)
        hex_val = "{:02X}{:02X}{:02X}".format(bg_color[0], bg_color[1], bg_color[2])
        fill_xml = (
            '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{hex_val}"/>'
            '</a:solidFill>'
        )
        tcPr.insert(0, etree.fromstring(fill_xml))


# -- Per-layout renderers ------------------------------------------------------

def render_title_slide(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["main_title"],   accent, content.get("main_title", ""))
    add_textbox(slide, boxes["sub_title"],    accent, content.get("sub_title", ""))
    date_co = content.get("date", "") + "  |  " + content.get("company", "")
    add_textbox(slide, boxes["bottom_info"],  accent, date_co.strip("  |  "))


def render_toc(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["toc_title"], accent, "CONTENTS")
    items  = content.get("toc_items", [])
    roman  = ["I.", "II.", "III.", "IV.", "V.", "VI.", "VII.", "VIII."]
    numbered = []
    for i, item in enumerate(items):
        prefix = roman[i] if i < len(roman) else f"{i+1}."
        numbered.append(f"{prefix:<5}{item}")
    add_textbox(slide, boxes["toc_items"], accent, "\n".join(numbered))


def render_section_divider(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["section_number"], accent, content.get("section_number", "Section"))
    add_textbox(slide, boxes["section_title"],  accent, content.get("section_title", ""))


def render_problem_vs_solution(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["top_message"],  accent, content.get("top_message", ""))
    add_textbox(slide, boxes["left_header"],  accent, content.get("left_title", ""))
    add_textbox(slide, boxes["left_body"],    accent, "\n".join(content.get("left_bullets", [])), bullet=True)
    add_textbox(slide, boxes["right_header"], accent, content.get("right_title", ""))
    add_textbox(slide, boxes["right_body"],   accent, "\n".join(content.get("right_bullets", [])), bullet=True)


def render_vendor_profile(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["top_message"], accent, content.get("top_message", ""))
    td      = content.get("table_data", {})
    headers = td.get("headers", [])
    rows    = td.get("rows", [])
    if headers and "table" in pos_def:
        add_table(slide, pos_def["table"], accent, headers, rows)


def render_flow_chart(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["top_message"], accent, content.get("top_message", ""))
    steps    = content.get("steps", [])
    step_ids = [
        ("step1_header", "step1_body"),
        ("step2_header", "step2_body"),
        ("step3_header", "step3_body"),
    ]
    for i, (hid, bid) in enumerate(step_ids):
        step = steps[i] if i < len(steps) else {}
        add_textbox(slide, boxes[hid], accent, step.get("header", f"STEP {i+1}"))
        add_textbox(slide, boxes[bid], accent, step.get("body", ""))


def render_comparison_benchmark(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["top_message"],  accent, content.get("top_message", ""))
    lb = content.get("left_block", {})
    rb = content.get("right_block", {})
    add_textbox(slide, boxes["left_header"],  accent, lb.get("title", ""))
    add_textbox(slide, boxes["left_body"],    accent, "\n".join(lb.get("bullets", [])), bullet=True)
    add_textbox(slide, boxes["right_header"], accent, rb.get("title", ""))
    add_textbox(slide, boxes["right_body"],   accent, "\n".join(rb.get("bullets", [])), bullet=True)
    td      = content.get("table_data", {})
    headers = td.get("headers", [])
    rows    = td.get("rows", [])
    if headers and "table" in pos_def:
        add_table(slide, pos_def["table"], accent, headers, rows)


def render_n_column_cards(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["top_message"], accent, content.get("top_message", ""))
    cards    = content.get("cards", [])
    card_ids = [
        ("card1_header", "card1_body"),
        ("card2_header", "card2_body"),
        ("card3_header", "card3_body"),
    ]
    for i, (hid, bid) in enumerate(card_ids):
        card = cards[i] if i < len(cards) else {}
        add_textbox(slide, boxes[hid], accent, card.get("header", ""))
        add_textbox(slide, boxes[bid], accent, card.get("body", ""))


def render_curriculum_table(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["slide_title"],  accent, content.get("top_message", ""))
    add_textbox(slide, boxes["course_goal"],  accent, content.get("course_goal", ""))
    td           = content.get("table_data", {})
    headers      = td.get("headers", [])
    rows         = td.get("rows", [])
    section_rows = td.get("section_rows", [])
    if section_rows:
        sr_map   = {sr["idx"]: sr["label"] for sr in section_rows}
        all_rows = []
        for idx, row in enumerate(rows):
            if idx in sr_map:
                sect_row = [sr_map[idx]] + [""] * (len(headers) - 1)
                all_rows.append(sect_row)
            all_rows.append(row)
        rows = all_rows
    if headers and "table" in pos_def:
        add_table(slide, pos_def["table"], accent, headers, rows)


def render_evaluation_metric(slide, pos_def, accent, content):
    boxes     = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["top_message"], accent, content.get("top_message", ""))
    cards     = content.get("cards", [])
    level_ids = [
        ("level1_header", "level1_body"),
        ("level2_header", "level2_body"),
        ("level3_header", "level3_body"),
        ("level4_header", "level4_body"),
    ]
    for i, (hid, bid) in enumerate(level_ids):
        card = cards[i] if i < len(cards) else {}
        add_textbox(slide, boxes[hid], accent, card.get("header", ""))
        add_textbox(slide, boxes[bid], accent, card.get("body", ""))
    if "footnote" in boxes:
        add_textbox(slide, boxes["footnote"], accent, content.get("footnote", ""))


def render_closing_slide(slide, pos_def, accent, content):
    boxes = {b["id"]: b for b in pos_def["boxes"]}
    add_textbox(slide, boxes["main_message"], accent, content.get("main_title", ""))
    add_textbox(slide, boxes["contact_info"], accent, content.get("contact_info", ""))


# -- Layout router -------------------------------------------------------------

LAYOUT_RENDERERS = {
    "TITLE_SLIDE":          render_title_slide,
    "TABLE_OF_CONTENTS":    render_toc,
    "SECTION_DIVIDER":      render_section_divider,
    "PROBLEM_VS_SOLUTION":  render_problem_vs_solution,
    "VENDOR_PROFILE":       render_vendor_profile,
    "FLOW_CHART":           render_flow_chart,
    "COMPARISON_BENCHMARK": render_comparison_benchmark,
    "N_COLUMN_CARDS":       render_n_column_cards,
    "CURRICULUM_TABLE":     render_curriculum_table,
    "EVALUATION_METRIC":    render_evaluation_metric,
    "CLOSING_SLIDE":        render_closing_slide,
}


# -- Main builder --------------------------------------------------------------

def build_proposal(proposal_json_path, output_path):
    with open(proposal_json_path, encoding="utf-8") as f:
        spec = json.load(f)

    with open(POSITIONS_PATH, encoding="utf-8") as f:
        positions = json.load(f)

    # Accent color from design_system
    design_sys  = spec.get("design_system", {})
    palette     = design_sys.get("color_palette", {})
    accent_hex  = palette.get("accent", "#1B5EA3")
    accent_color = hex_to_rgb(accent_hex)

    print("[INFO] Accent: " + accent_hex)
    print("[INFO] Slides: " + str(len(spec.get("slides", []))))

    prs = Presentation()
    prs.slide_width  = Cm(SLIDE_W_CM)
    prs.slide_height = Cm(SLIDE_H_CM)

    for slide_spec in spec.get("slides", []):
        layout_type = slide_spec.get("layout_type", "")
        content     = slide_spec.get("content", {})
        slide_num   = slide_spec.get("slide_number", "?")

        if layout_type not in positions:
            print("[WARN] Slide " + str(slide_num) + ": layout '" + layout_type + "' not in positions.json -- skip")
            continue

        if layout_type not in LAYOUT_RENDERERS:
            print("[WARN] Slide " + str(slide_num) + ": no renderer for '" + layout_type + "' -- skip")
            continue

        slide    = add_blank_slide(prs)
        pos_def  = positions[layout_type]
        renderer = LAYOUT_RENDERERS[layout_type]
        renderer(slide, pos_def, accent_color, content)
        print("[OK]  Slide " + str(slide_num) + ": " + layout_type)

    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    print("\n[DONE] Saved: " + output_path)


# -- Entry point ---------------------------------------------------------------

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python renderer/pptx_builder.py <proposal_json> <output_pptx>")
        print("")
        print("Example:")
        print("  python renderer/pptx_builder.py output/dongdaemun_proposal.json output/dongdaemun_v3.pptx")
        sys.exit(1)

    json_path   = sys.argv[1]
    output_path = sys.argv[2]

    if not os.path.exists(json_path):
        print("[ERROR] JSON not found: " + json_path)
        sys.exit(1)

    build_proposal(json_path, output_path)
