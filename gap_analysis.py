import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

def emu_in(e):
    return Emu(e or 0).inches

def pt(e):
    return (e or 0) / 12700.0

print("=== Gap between title box bottom and green header bar top ===")
print("  (title clips if its box is too short AND no gap to green bar)")
print()

for i, slide in enumerate(prs.slides, 1):
    if i < 3 or i > 9:
        continue

    # Find title box (top < 0.6")
    title_shape = None
    for s in slide.shapes:
        if s.has_text_frame and emu_in(s.top) < 0.6:
            title_shape = s
            break

    # Find first green header bar (top ~ 1.3-1.4")
    green_bars = []
    for s in slide.shapes:
        if s.has_text_frame and 1.0 <= emu_in(s.top) <= 1.5:
            green_bars.append(s)

    if title_shape is None:
        print(f"Slide {i}: no title shape found")
        continue

    t_top = emu_in(title_shape.top)
    t_h = emu_in(title_shape.height)
    t_bottom = t_top + t_h
    title_text = title_shape.text_frame.text.replace('\n', ' | ')

    # Get font size
    max_pt = 0
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                v = pt(run.font.size)
                if v > max_pt:
                    max_pt = v

    # Estimated actual render height of two-line title
    # At 30pt with default 100% line spacing, each line = 30/72 = 0.417"
    # But PowerPoint's default line spacing for Korean text can be tighter
    # Let's use 1.15x (115% is common default)
    line_h = max_pt / 72.0
    inset_top = emu_in(45720)  # 0.05"
    inset_bot = emu_in(45720)  # 0.05"
    render_h_tight = line_h * 2 + inset_top + inset_bot  # tight estimate
    render_h_normal = max_pt * 1.2 / 72 * 2 + inset_top + inset_bot  # 120% line spacing

    print(f"Slide {i}:")
    print(f"  Title: '{title_text[:60]}'")
    print(f"  Font: {max_pt:.1f}pt, box: top={t_top:.3f}\" h={t_h:.3f}\" bottom={t_bottom:.3f}\"")
    print(f"  Render est (tight/normal): {render_h_tight:.3f}\"/{ render_h_normal:.3f}\" vs box {t_h:.3f}\"")

    if green_bars:
        bar_top = min(emu_in(s.top) for s in green_bars)
        gap = bar_top - t_bottom
        print(f"  Green bar top: {bar_top:.3f}\"  gap from title bottom: {gap:.3f}\"")

        # The critical question: does second title line render below title box bottom?
        # With spAutoFit=yes, PowerPoint will expand the box if text overflows
        # BUT the visual issue is when title text overflows BEHIND the green bar
        second_line_bottom = t_top + render_h_normal
        if second_line_bottom > bar_top:
            print(f"  *** OVERLAP RISK: second line renders to {second_line_bottom:.3f}\" but green bar starts at {bar_top:.3f}\" ***")
        else:
            print(f"  OK: second line ends {second_line_bottom:.3f}\" before bar at {bar_top:.3f}\" (gap={bar_top - second_line_bottom:.3f}\")")
    else:
        print(f"  No green bar found")
    print()
