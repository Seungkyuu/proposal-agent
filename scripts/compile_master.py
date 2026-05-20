"""
compile_master.py
─────────────────────────────────────────────────────────────────────
slide_catalog.json을 읽어 각 소스 PPTX에서 추천 슬라이드를 복제하여
하나의 alpaco_master.pptx 마스터 덱을 생성한다.

실행:
    pip install python-pptx
    python scripts/compile_master.py

출력:
    templates/master/alpaco_master.pptx   ← 마스터 덱 (레이아웃 모음)
    templates/master/slide_index.json     ← layout_type → master 슬라이드 번호 매핑
"""

import json
import os
import copy
from datetime import date
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ─── 경로 설정 ────────────────────────────────────────────────────
BASE_DIR     = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CATALOG_PATH = os.path.join(BASE_DIR, 'templates', 'slide_catalog.json')
SOURCE_DIR   = os.path.join(BASE_DIR, 'templates', 'source_decks')
MASTER_DIR   = os.path.join(BASE_DIR, 'templates', 'master')
MASTER_PPTX  = os.path.join(MASTER_DIR, 'alpaco_master.pptx')
INDEX_PATH   = os.path.join(MASTER_DIR, 'slide_index.json')

os.makedirs(MASTER_DIR, exist_ok=True)


# ─── 슬라이드 크로스파일 복제 헬퍼 ──────────────────────────────────
def _copy_slide(src_prs, src_slide_idx, dst_prs):
    """
    src_prs의 src_slide_idx(0-based) 슬라이드를
    dst_prs 끝에 복제하여 추가한다.
    python-pptx 공식 지원 외 XML 직접 조작 방식.
    """
    src_slide = src_prs.slides[src_slide_idx]

    # 슬라이드 레이아웃: 대상 프레젠테이션의 blank 레이아웃 사용
    slide_layout = dst_prs.slide_layouts[6]  # blank
    new_slide = dst_prs.slides.add_slide(slide_layout)

    # 기존 placeholder/shape 모두 제거
    sp_tree = new_slide.shapes._spTree
    for el in list(sp_tree):
        sp_tree.remove(el)

    # 원본 슬라이드 XML 요소를 깊은 복사하여 이식
    src_sp_tree = src_slide.shapes._spTree
    for el in src_sp_tree:
        sp_tree.append(copy.deepcopy(el))

    # 슬라이드 크기 동기화 (원본 → 대상)
    dst_prs.slide_width  = src_prs.slide_width
    dst_prs.slide_height = src_prs.slide_height

    return new_slide


# ─── 메인 실행 ────────────────────────────────────────────────────
def main():
    print("\n─── 마스터 덱 컴파일 시작 ───\n")

    with open(CATALOG_PATH, encoding='utf-8') as f:
        catalog = json.load(f)

    layouts = catalog['layouts']

    # 레이아웃 타입 순서 고정 (제안서 실제 슬라이드 순서와 맞춤)
    LAYOUT_ORDER = [
        'TITLE_SLIDE',
        'TABLE_OF_CONTENTS',
        'SECTION_DIVIDER',
        'VENDOR_PROFILE',
        'PROBLEM_VS_SOLUTION',
        'FLOW_CHART',
        'COMPARISON_BENCHMARK',
        'CURRICULUM_TABLE',
        'N_COLUMN_CARDS',
        'EVALUATION_METRIC',
        'CLOSING_SLIDE',
    ]

    # 마스터 덱 프레젠테이션 초기화
    master_prs = Presentation()
    master_prs.slide_width  = Inches(13.33)
    master_prs.slide_height = Inches(7.5)

    slide_index = {
        "master_file": "alpaco_master.pptx",
        "version": "1.0",
        "generated_at": str(date.today()),
        "layouts": {}
    }

    master_slide_num = 0  # 0-based counter

    # 캐시: 같은 소스 파일을 여러 번 열지 않도록
    prs_cache = {}

    for layout_type in LAYOUT_ORDER:
        if layout_type not in layouts:
            print(f"  ⚠  {layout_type}: catalog에 없음 — 건너뜀")
            continue

        rec  = layouts[layout_type]['recommended']
        fname = rec['file']
        sidx  = rec['slide_idx'] - 1  # 0-based

        fpath = os.path.join(SOURCE_DIR, fname)
        if not os.path.exists(fpath):
            print(f"  ⚠  {layout_type}: 파일 없음 — {fname}")
            continue

        # 소스 파일 로드 (캐시)
        if fname not in prs_cache:
            prs_cache[fname] = Presentation(fpath)
        src_prs = prs_cache[fname]

        if sidx >= len(src_prs.slides):
            print(f"  ⚠  {layout_type}: 슬라이드 인덱스 초과 ({sidx+1}/{len(src_prs.slides)})")
            continue

        # 복제
        _copy_slide(src_prs, sidx, master_prs)
        master_slide_num += 1

        slide_index['layouts'][layout_type] = {
            "slide_idx": master_slide_num,          # 1-based (마스터 덱 내 순번)
            "description": layouts[layout_type]['description'],
            "source_file": fname,
            "source_slide_idx": rec['slide_idx'],
            "flags": rec.get('flags', ''),
            # 렌더러가 텍스트 교체 시 참고하는 힌트
            "content_hints": _get_content_hints(layout_type),
        }

        print(f"  ✓  [{master_slide_num:2d}] {layout_type:<25}  ←  {fname[:45]} S{rec['slide_idx']:02d}")

    # 저장
    master_prs.save(MASTER_PPTX)

    with open(INDEX_PATH, 'w', encoding='utf-8') as f:
        json.dump(slide_index, f, ensure_ascii=False, indent=2)

    print(f"\n✅ 마스터 덱 저장: {MASTER_PPTX}")
    print(f"✅ 슬라이드 인덱스 저장: {INDEX_PATH}")
    print(f"   총 슬라이드: {master_slide_num}장")
    print("\n📋 다음 단계:")
    print("   1. PowerPoint로 alpaco_master.pptx 열어 레이아웃 시각 확인")
    print("   2. 마음에 안 드는 슬라이드 → slide_catalog.json alternatives 중 교체 후 재실행")
    print("   3. 확인 완료 후 python renderer/pptx_builder.py 로 제안서 생성\n")


def _get_content_hints(layout_type):
    """
    pptx_builder.py가 텍스트 교체 시 참고하는 구조 힌트.
    실제 placeholder 위치는 마스터 덱 분석 후 보완 가능.
    """
    hints = {
        'TITLE_SLIDE': {
            'fields': ['main_title', 'sub_title', 'date', 'company'],
            'strategy': 'text_replace_by_order',
        },
        'TABLE_OF_CONTENTS': {
            'fields': ['toc_items'],
            'strategy': 'bullet_replace',
        },
        'SECTION_DIVIDER': {
            'fields': ['section_title', 'section_subtitle'],
            'strategy': 'text_replace_by_order',
        },
        'VENDOR_PROFILE': {
            'fields': ['top_message', 'table_data'],
            'strategy': 'text_and_table',
            'table_cols': ['프로젝트명', '고객사', '규모', '만족도'],
        },
        'PROBLEM_VS_SOLUTION': {
            'fields': ['top_message', 'left_title', 'left_bullets', 'right_title', 'right_bullets'],
            'strategy': 'two_column_replace',
        },
        'FLOW_CHART': {
            'fields': ['top_message', 'steps'],
            'strategy': 'card_replace',
            'card_count': 3,
        },
        'COMPARISON_BENCHMARK': {
            'fields': ['top_message', 'left_title', 'left_bullets', 'right_title', 'right_bullets', 'alignment_score'],
            'strategy': 'two_column_replace',
        },
        'CURRICULUM_TABLE': {
            'fields': ['top_message', 'table_data'],
            'strategy': 'text_and_table',
            'table_cols': ['차시', '세부 주제', '핵심 학습 내용', '실습 도구 및 액티비티', '최종 산출물'],
            'supports_section_rows': True,
        },
        'N_COLUMN_CARDS': {
            'fields': ['top_message', 'cards'],
            'strategy': 'card_replace',
            'card_count': 3,
        },
        'EVALUATION_METRIC': {
            'fields': ['top_message', 'cards'],
            'strategy': 'card_replace',
            'card_count': 4,
        },
        'CLOSING_SLIDE': {
            'fields': ['main_title', 'contact_info'],
            'strategy': 'text_replace_by_order',
        },
    }
    return hints.get(layout_type, {'strategy': 'text_replace_by_order'})


if __name__ == '__main__':
    main()
