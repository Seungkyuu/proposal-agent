import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
import os

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

slide_w = prs.slide_width
slide_h = prs.slide_height
SLIDE_W_IN = Emu(slide_w).inches
SLIDE_H_IN = Emu(slide_h).inches

def emu_in(e):
    if e is None:
        return 0.0
    return Emu(e).inches

for i, slide in enumerate(prs.slides, 1):
    print(f"=== Slide {i} ===")
    shapes = list(slide.shapes)

    for shape in shapes:
        l = emu_in(shape.left)
        t = emu_in(shape.top)
        w = emu_in(shape.width)
        h = emu_in(shape.height)
        r = l + w
        b = t + h

        flags = []
        if r > SLIDE_W_IN + 0.01:
            flags.append(f"OVERFLOW_RIGHT by {r-SLIDE_W_IN:.3f}\"")
        if b > SLIDE_H_IN + 0.01:
            flags.append(f"OVERFLOW_BOTTOM by {b-SLIDE_H_IN:.3f}\"")
        if l < -0.01:
            flags.append(f"OFF_LEFT by {-l:.3f}\"")
        if t < -0.01:
            flags.append(f"OFF_TOP by {-t:.3f}\"")

        text = ""
        font_info = []
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text[:100].replace('\n', ' | ')
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_info.append(f"{Pt(run.font.size).pt:.0f}pt")
            # Word wrap setting
            wrap = tf.word_wrap
            # auto_size
            auto = tf.auto_size

        name = shape.shape_type
        print(f"  [{shape.name}] pos=({l:.2f}\",{t:.2f}\") size=({w:.2f}\"x{h:.2f}\")", end="")
        if text:
            print(f"\n    text: '{text[:70]}'", end="")
        if font_info:
            print(f"\n    fonts: {', '.join(set(font_info))}", end="")
        if flags:
            for f in flags:
                print(f"\n    *** {f} ***", end="")
        print()

    # Check title text box specifically for two-line titles
    title_shapes = [s for s in shapes if s.has_text_frame and
                    emu_in(s.top) < 0.8 and emu_in(s.height) < 1.2]
    for ts in title_shapes:
        tf = ts.text_frame
        raw = tf.text
        lines = raw.split('\n')
        # Check if title has pipe separator suggesting two-line layout
        if '|' in raw:
            parts = [p.strip() for p in raw.split('|')]
            print(f"  TITLE HAS {len(parts)} PARTS (pipe-separated): may wrap")
            for j, p in enumerate(parts):
                print(f"    Part {j+1}: '{p[:60]}'")
        # Approximate character-based overflow check
        # At ~20pt bold, ~0.12" per char on 13.33" wide box
        for para in tf.paragraphs:
            for run in para.runs:
                fs = run.font.size
                if fs:
                    pt = Pt(fs).pt
                    chars = len(run.text)
                    # rough: chars * pt * 0.006 inches per char per pt
                    est_width = chars * pt * 0.006
                    box_w = emu_in(ts.width)
                    if est_width > box_w:
                        print(f"  POSSIBLE TEXT OVERFLOW in run: '{run.text[:50]}' est_width={est_width:.2f}\" box={box_w:.2f}\"")

    # Check for overlapping shapes
    shape_rects = []
    for s in shapes:
        l = emu_in(s.left)
        t = emu_in(s.top)
        r = l + emu_in(s.width)
        b = t + emu_in(s.height)
        shape_rects.append((s.name, l, t, r, b))

    for j in range(len(shape_rects)):
        for k in range(j+1, len(shape_rects)):
            n1, l1, t1, r1, b1 = shape_rects[j]
            n2, l2, t2, r2, b2 = shape_rects[k]
            # Check overlap
            ox = min(r1, r2) - max(l1, l2)
            oy = min(b1, b2) - max(t1, t2)
            if ox > 0.05 and oy > 0.05:
                print(f"  OVERLAP: '{n1}' and '{n2}' overlap by ({ox:.2f}\"x{oy:.2f}\")")

    print()
