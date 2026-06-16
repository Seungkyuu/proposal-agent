import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu, Pt
from lxml import etree

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

def emu_to_pt(v):
    return v / 12700.0 if v else 0

# Check line spacing settings on title paragraphs
print("=== Title paragraph line spacing settings ===")
NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t_in = Emu(shape.top or 0).inches
        h_in = Emu(shape.height or 0).inches
        if t_in > 0.6:
            continue
        tf = shape.text_frame
        print(f"Slide {i}: top={t_in:.3f}\" height={h_in:.3f}\"")
        # Check txBody attributes
        txBody = tf._txBody
        # bodyPr
        bodyPr = txBody.find(f'{{{NS}}}bodyPr')
        if bodyPr is not None:
            wrap = bodyPr.get('wrap', 'square')
            autofit = bodyPr.find(f'{{{NS}}}normAutofit')
            noAutofit = bodyPr.find(f'{{{NS}}}noAutofit')
            spAutoFit = bodyPr.find(f'{{{NS}}}spAutoFit')
            print(f"  bodyPr: wrap={wrap}, normAutofit={'yes' if autofit is not None else 'no'}, noAutofit={'yes' if noAutofit is not None else 'no'}, spAutoFit={'yes' if spAutoFit is not None else 'no'}")
            ins_t = bodyPr.get('insTfm', bodyPr.get('insT', '45720'))  # top inset
            ins_b = bodyPr.get('insB', '45720')  # bottom inset
            # inset in EMU: default 45720 = 0.05" top/bottom, 91440 = 0.1" left/right
            print(f"  insets: top={bodyPr.get('insT','45720')} bottom={bodyPr.get('insB','45720')}")

        for pi, para in enumerate(tf.paragraphs):
            ptext = para.text[:40]
            # Check paragraph properties
            pPr = para._p.find(f'{{{NS}}}pPr')
            if pPr is not None:
                # Line spacing
                lnSpc = pPr.find(f'{{{NS}}}lnSpc')
                spcBef = pPr.find(f'{{{NS}}}spcBef')
                spcAft = pPr.find(f'{{{NS}}}spcAft')
                ls_info = "default(100%)"
                if lnSpc is not None:
                    spcPct = lnSpc.find(f'{{{NS}}}spcPct')
                    spcPts = lnSpc.find(f'{{{NS}}}spcPts')
                    if spcPct is not None:
                        ls_info = f"spcPct={spcPct.get('val')}%"
                    elif spcPts is not None:
                        ls_info = f"spcPts={spcPts.get('val')}"
                sb_info = "0"
                if spcBef is not None:
                    spcPts = spcBef.find(f'{{{NS}}}spcPts')
                    spcPct = spcBef.find(f'{{{NS}}}spcPct')
                    if spcPts is not None:
                        sb_info = f"{int(spcPts.get('val',0))/100:.1f}pt"
                    elif spcPct is not None:
                        sb_info = f"{spcPct.get('val')}%"
                print(f"  para[{pi}]: '{ptext}' lineSpacing={ls_info} spaceBefore={sb_info}")
        print()

print()
print("=== Green header bar shapes (if any rectangles near top on slides 3-9) ===")
for i, slide in enumerate(prs.slides, 1):
    if i < 3 or i > 9:
        continue
    for shape in slide.shapes:
        t_in = Emu(shape.top or 0).inches
        h_in = Emu(shape.height or 0).inches
        # Check for solid rectangles / auto shapes near top
        if t_in < 1.5 and not shape.has_text_frame:
            print(f"  Slide {i}: shape '{shape.name}' type={shape.shape_type} top={t_in:.3f}\" height={h_in:.3f}\"")
        elif t_in < 1.5 and shape.has_text_frame and shape.text_frame.text == '':
            print(f"  Slide {i}: empty text box '{shape.name}' top={t_in:.3f}\" height={h_in:.3f}\"")
