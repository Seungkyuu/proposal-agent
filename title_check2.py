import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu, Pt

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

print("=== All title shapes (top < 0.6\") with font sizes ===")
print()

for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t_in = Emu(shape.top or 0).inches
        h_in = Emu(shape.height or 0).inches
        if t_in > 0.6:
            continue
        tf = shape.text_frame
        text = tf.text
        print(f"Slide {i}: top={t_in:.3f}\" height={h_in:.3f}\" | text='{text[:80]}'")
        # Get actual font sizes from paragraphs
        for pi, para in enumerate(tf.paragraphs):
            ptext = para.text[:60]
            pfont = None
            for run in para.runs:
                if run.font.size:
                    pfont = Pt(run.font.size).pt
                    break
            if pfont is None and para.runs:
                pfont = "inherit"
            print(f"  para[{pi}]: '{ptext}' font={pfont}")
        print()

print()
print("=== Slide 1 detailed layout ===")
slide1 = prs.slides[0]
for shape in sorted(slide1.shapes, key=lambda s: Emu(s.top or 0).inches):
    t = Emu(shape.top or 0).inches
    h = Emu(shape.height or 0).inches
    b = t + h
    if shape.has_text_frame:
        text = shape.text_frame.text[:80].replace('\n', ' / ')
        # Get font sizes
        sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(f"{Pt(run.font.size).pt:.0f}pt")
        print(f"  top={t:.3f}\" h={h:.3f}\" bottom={b:.3f}\" fonts={set(sizes)} | '{text}'")

print()
print("=== Slide 2 detailed layout ===")
slide2 = prs.slides[1]
for shape in sorted(slide2.shapes, key=lambda s: Emu(s.top or 0).inches):
    t = Emu(shape.top or 0).inches
    h = Emu(shape.height or 0).inches
    b = t + h
    if shape.has_text_frame:
        text = shape.text_frame.text[:80].replace('\n', ' / ')
        sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(f"{Pt(run.font.size).pt:.0f}pt")
        print(f"  top={t:.3f}\" h={h:.3f}\" bottom={b:.3f}\" fonts={set(sizes)} | '{text}'")
