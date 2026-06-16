import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu

pptx_path = r'C:\Users\smvo0\proposal-agent\output\dongdaemun_v3.pptx'
prs = Presentation(pptx_path)

def emu_in(e):
    return Emu(e or 0).inches

def pt(e):
    return (e or 0) / 12700.0

print("=== Slide 7 layout detail ===")
slide7 = prs.slides[6]
for shape in sorted(slide7.shapes, key=lambda s: emu_in(s.top)):
    t = emu_in(shape.top)
    h = emu_in(shape.height)
    b = t + h
    w = emu_in(shape.width)
    if shape.has_text_frame:
        text = shape.text_frame.text[:60].replace('\n', ' / ')
        fonts = set()
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    fonts.add(f"{pt(run.font.size):.0f}pt")
        print(f"  top={t:.3f}\" h={h:.3f}\" w={w:.2f}\" bottom={b:.3f}\" fonts={fonts} | '{text}'")
    else:
        print(f"  top={t:.3f}\" h={h:.3f}\" w={w:.2f}\" bottom={b:.3f}\" [{shape.shape_type}] '{shape.name}'")

print()
print("=== Slide 7: title box vs table top ===")
title = None
table = None
for shape in slide7.shapes:
    t = emu_in(shape.top)
    if shape.has_text_frame and t < 0.6:
        title = shape
    if shape.shape_type == 19:
        table = shape

if title and table:
    t_bot = emu_in(title.top) + emu_in(title.height)
    tbl_top = emu_in(table.top)
    gap = tbl_top - t_bot
    print(f"  Title bottom={t_bot:.3f}\" | Table top={tbl_top:.3f}\" | Gap={gap:.3f}\"")
    # Font size from title
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                font_pt = pt(run.font.size)
                line_h = font_pt / 72.0
                render_h = font_pt * 1.2 / 72 * 2 + 0.1  # two lines + insets
                print(f"  Font={font_pt:.0f}pt, 2-line render est={render_h:.3f}\" vs box={emu_in(title.height):.3f}\"")
                # Does expanded title overlap empty textbox or table?
                actual_bottom = emu_in(title.top) + render_h
                print(f"  Rendered bottom est={actual_bottom:.3f}\"")
                break

print()
print("=== Slide 10 layout ===")
slide10 = prs.slides[9]
SLIDE_W = Emu(prs.slide_width).inches
SLIDE_H = Emu(prs.slide_height).inches
for shape in sorted(slide10.shapes, key=lambda s: emu_in(s.top)):
    t = emu_in(shape.top)
    h = emu_in(shape.height)
    b = t + h
    w = emu_in(shape.width)
    l = emu_in(shape.left)
    if shape.has_text_frame:
        text = shape.text_frame.text[:80].replace('\n', ' / ')
        print(f"  left={l:.2f}\" top={t:.3f}\" w={w:.2f}\" h={h:.3f}\" bottom={b:.3f}\" | '{text}'")
    else:
        print(f"  left={l:.2f}\" top={t:.3f}\" w={w:.2f}\" h={h:.3f}\" bottom={b:.3f}\" [{shape.name}]")

print()
print(f"Slide size: {SLIDE_W:.2f}\" x {SLIDE_H:.2f}\"")

print()
print("=== Slide 1: gap checks ===")
slide1 = prs.slides[0]
shapes_s1 = sorted(slide1.shapes, key=lambda s: emu_in(s.top))
for j, s in enumerate(shapes_s1):
    t = emu_in(s.top)
    h = emu_in(s.height)
    b = t + h
    text = s.text_frame.text[:50].replace('\n', ' / ') if s.has_text_frame else ""
    print(f"  top={t:.3f}\" bottom={b:.3f}\" | '{text}'")
    if j < len(shapes_s1) - 1:
        next_t = emu_in(shapes_s1[j+1].top)
        gap = next_t - b
        print(f"    -> gap to next: {gap:.3f}\"")
